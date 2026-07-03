"""
Universal Document Reader
Reads text from PDF, DOCX, PPTX, and other document formats.

Usage: python read_document.py <file_path> [--pages 1-5] [--max-chars 50000]
"""
import sys
import argparse
from pathlib import Path


def read_pdf(file_path: str, pages: str = None, max_chars: int = None) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        total_pages = len(pdf.pages)
        page_indices = parse_page_range(pages, total_pages) if pages else range(total_pages)
        for i in page_indices:
            if i < total_pages:
                page = pdf.pages[i]
                page_text = page.extract_text() or ""
                text_parts.append(f"--- Page {i + 1} ---\n{page_text}")
    return "\n\n".join(text_parts)


def read_docx(file_path: str, max_chars: int = None) -> str:
    from docx import Document
    doc = Document(file_path)
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            style = para.style.name if para.style else ""
            prefix = ""
            if "Heading" in style:
                level = style.replace("Heading", "").strip() or "1"
                prefix = "#" * int(level) + " "
            text_parts.append(f"{prefix}{para.text}")
    for table in doc.tables:
        text_parts.append("\n[Table]")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            text_parts.append(" | ".join(cells))
    return "\n\n".join(text_parts)


def read_pptx(file_path: str, max_chars: int = None) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text)
            if shape.has_table:
                slide_text.append("[Table]")
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_text.append(" | ".join(cells))
        if len(slide_text) > 1:
            text_parts.append("\n".join(slide_text))
    return "\n\n".join(text_parts)


def parse_page_range(pages_str: str, total: int) -> range:
    if "-" in pages_str:
        start, end = pages_str.split("-", 1)
        return range(int(start) - 1, min(int(end), total))
    elif "," in pages_str:
        return [int(p) - 1 for p in pages_str.split(",") if int(p) <= total]
    else:
        p = int(pages_str) - 1
        return range(p, min(p + 1, total))


def main():
    parser = argparse.ArgumentParser(description="Read text from document files")
    parser.add_argument("file_path", help="Path to the document file")
    parser.add_argument("--pages", help="Page range (e.g., 1-5, 1,3,7)", default=None)
    parser.add_argument("--max-chars", type=int, default=100000, help="Max characters to return")
    args = parser.parse_args()

    path = Path(args.file_path)
    if not path.exists():
        print(f"Error: File not found: {args.file_path}", file=sys.stderr)
        sys.exit(1)

    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            text = read_pdf(args.file_path, args.pages, args.max_chars)
        elif ext == ".docx":
            text = read_docx(args.file_path, args.max_chars)
        elif ext == ".pptx":
            text = read_pptx(args.file_path, args.max_chars)
        else:
            print(f"Error: Unsupported file type '{ext}'. Supported: .pdf, .docx, .pptx", file=sys.stderr)
            sys.exit(1)

        if args.max_chars and len(text) > args.max_chars:
            text = text[:args.max_chars] + f"\n\n[Truncated at {args.max_chars} chars]"

        print(text if text else "[No text content found in document]")

    except Exception as e:
        print(f"Error reading file: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
