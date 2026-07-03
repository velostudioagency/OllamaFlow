import os
import re
import subprocess
import sys
import time as _time
from pathlib import Path
from typing import Any, Dict

from app.core.config import WORKSPACE_DIR


def read_file(file_path: str, **kwargs) -> str:
    try:
        path = Path(file_path)
        if not path.is_absolute():
            path = Path(WORKSPACE_DIR) / file_path
        if not path.exists():
            return f"Error: File not found: {file_path}"
        suffix = path.suffix.lower()
        if suffix == ".txt" or suffix == ".md" or suffix == ".py" or suffix == ".json" or suffix == ".csv":
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            if len(content) > 10000:
                content = content[:10000] + "\n... [truncated]"
            return content
        elif suffix == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(str(path)) as pdf:
                    text_parts = []
                    for i, page in enumerate(pdf.pages[:30]):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            text_parts.append(f"--- Page {i + 1} ---\n{page_text}")
                    return "\n\n".join(text_parts)[:10000]
            except ImportError:
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(str(path))
                    text_parts = []
                    for page in reader.pages[:20]:
                        text_parts.append(page.extract_text() or "")
                    return "\n".join(text_parts)[:10000]
                except ImportError:
                    return "Error: Neither pdfplumber nor pypdf installed"
        elif suffix == ".docx":
            try:
                from docx import Document
                doc = Document(str(path))
                text_parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                for table in doc.tables:
                    text_parts.append("[Table]")
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        text_parts.append(" | ".join(cells))
                return "\n".join(text_parts)[:10000]
            except ImportError:
                return "Error: python-docx not installed"
        elif suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                text_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = [f"--- Slide {slide_num} ---"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    slide_text.append(para.text)
                        if shape.has_table:
                            slide_text.append("[Table]")
                            for row in shape.table.rows:
                                cells = [cell.text.strip() for cell in row.cells]
                                slide_text.append(" | ".join(cells))
                    if len(slide_text) > 1:
                        text_parts.append("\n".join(slide_text))
                return "\n\n".join(text_parts)[:10000]
            except ImportError:
                return "Error: python-pptx not installed"
        elif suffix in [".xlsx", ".xls"]:
            try:
                import pandas as pd
                df = pd.read_excel(str(path))
                return df.to_string()[:10000]
            except ImportError:
                return "Error: openpyxl not installed"
        else:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()[:10000]
    except Exception as e:
        return f"Error reading file: {str(e)}"


def write_file(file_path: str, content: str, **kwargs) -> str:
    try:
        path = Path(file_path)
        if not path.is_absolute():
            path = Path(WORKSPACE_DIR) / file_path
        path.parent.mkdir(parents=True, exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return f"Successfully wrote {len(content)} bytes to {path}"
    except Exception as e:
        return f"Error writing file: {str(e)}"


def list_directory(path: str = "", pattern: str = "", **kwargs) -> str:
    try:
        dir_path = Path(path) if path else Path(WORKSPACE_DIR)
        if not dir_path.is_absolute():
            dir_path = Path(WORKSPACE_DIR) / path
        if not dir_path.exists():
            return f"Error: Directory not found: {path}"
        if not dir_path.is_dir():
            return f"Error: Not a directory: {path}"
        entries = []
        glob_pattern = pattern if pattern else "*"
        for item in sorted(dir_path.glob(glob_pattern)):
            if item.name.startswith("."):
                continue
            try:
                size = item.stat().st_size
                if size < 1024:
                    size_str = f"{size}B"
                elif size < 1024 * 1024:
                    size_str = f"{size / 1024:.1f}KB"
                else:
                    size_str = f"{size / (1024 * 1024):.1f}MB"
            except Exception:
                size_str = "?"
            item_type = "DIR " if item.is_dir() else "FILE"
            entries.append(f"  [{item_type}] {item.name:40s} {size_str}")
        if not entries:
            return f"Directory is empty: {dir_path}"
        header = f"Contents of {dir_path}:\n"
        return header + "\n".join(entries)
    except Exception as e:
        return f"Error listing directory: {str(e)}"


def search_files(query: str, path: str = "", search_type: str = "name", **kwargs) -> str:
    try:
        search_dir = Path(path) if path else Path(WORKSPACE_DIR)
        if not search_dir.is_absolute():
            search_dir = Path(WORKSPACE_DIR) / path
        if not search_dir.exists():
            return f"Error: Directory not found: {path}"
        results = []
        if search_type == "name":
            for item in search_dir.rglob(f"*{query}*"):
                if item.name.startswith(".") or "node_modules" in str(item):
                    continue
                rel = item.relative_to(search_dir)
                results.append(str(rel))
                if len(results) >= 50:
                    break
        elif search_type == "content":
            pattern = re.compile(query, re.IGNORECASE)
            for item in search_dir.rglob("*"):
                if item.name.startswith(".") or "node_modules" in str(item):
                    continue
                if not item.is_file():
                    continue
                try:
                    text = item.read_text(encoding="utf-8", errors="replace")[:50000]
                    for i, line in enumerate(text.split("\n"), 1):
                        if pattern.search(line):
                            rel = item.relative_to(search_dir)
                            results.append(f"{rel}:{i}: {line.strip()[:120]}")
                            if len(results) >= 50:
                                break
                except Exception:
                    continue
                if len(results) >= 50:
                    break
        if not results:
            return f"No results found for '{query}' in {search_dir}"
        return f"Found {len(results)} results:\n" + "\n".join(results)
    except Exception as e:
        return f"Error searching files: {str(e)}"


def file_watcher(path: str = "", pattern: str = "*", action: str = "list_changes", since_minutes: int = 60, **kwargs) -> str:
    """Watch or list recent file changes in a directory."""
    try:
        watch_dir = Path(path) if path else Path(WORKSPACE_DIR)
        if not watch_dir.is_absolute():
            watch_dir = Path(WORKSPACE_DIR) / path
        if not watch_dir.exists():
            return f"Error: Directory not found: {path}"
        cutoff = _time.time() - (since_minutes * 60)
        changes = []
        for item in watch_dir.rglob(pattern):
            if item.name.startswith(".") or "node_modules" in str(item):
                continue
            try:
                mtime = item.stat().st_mtime
                if mtime >= cutoff:
                    age_s = _time.time() - mtime
                    if age_s < 60:
                        age = f"{age_s:.0f}s ago"
                    elif age_s < 3600:
                        age = f"{age_s/60:.0f}m ago"
                    else:
                        age = f"{age_s/3600:.0f}h ago"
                    changes.append({
                        "path": str(item.relative_to(watch_dir)),
                        "type": "dir" if item.is_dir() else "file",
                        "modified": age,
                        "size": item.stat().st_size
                    })
            except Exception:
                continue
        if not changes:
            return f"No changes in last {since_minutes} minutes in {watch_dir}"
        lines = [f"Changes in {watch_dir} (last {since_minutes} min):", ""]
        for c in sorted(changes, key=lambda x: x["path"]):
            size = c["size"]
            if size < 1024:
                size_str = f"{size}B"
            elif size < 1024 * 1024:
                size_str = f"{size/1024:.1f}KB"
            else:
                size_str = f"{size/(1024*1024):.1f}MB"
            lines.append(f"  [{c['type'].upper()}] {c['path']:40s} {size_str:10s} {c['modified']}")
        return "\n".join(lines)[:5000]
    except Exception as e:
        return f"File watcher error: {str(e)}"


def clipboard_copy(text: str = "", **kwargs) -> str:
    """Copy text to the system clipboard."""
    try:
        if sys.platform == "win32":
            process = subprocess.Popen(["clip"], stdin=subprocess.PIPE)
            process.communicate(text.encode("utf-16le"))
        elif sys.platform == "darwin":
            process = subprocess.Popen(["pbcopy"], stdin=subprocess.PIPE)
            process.communicate(text.encode("utf-8"))
        else:
            process = subprocess.Popen(["xclip", "-selection", "clipboard"], stdin=subprocess.PIPE)
            process.communicate(text.encode("utf-8"))
        return f"Copied {len(text)} characters to clipboard."
    except Exception as e:
        return f"Clipboard copy error: {str(e)}"


def clipboard_paste(**kwargs) -> str:
    """Read text from the system clipboard."""
    try:
        if sys.platform == "win32":
            result = subprocess.run(["powershell", "-command", "Get-Clipboard"],
                                    capture_output=True, text=True, timeout=5)
            return result.stdout.strip()
        elif sys.platform == "darwin":
            result = subprocess.run(["pbpaste"], capture_output=True, text=True, timeout=5)
            return result.stdout.strip()
        else:
            result = subprocess.run(["xclip", "-selection", "clipboard", "-o"],
                                    capture_output=True, text=True, timeout=5)
            return result.stdout.strip()
    except Exception as e:
        return f"Clipboard paste error: {str(e)}"
