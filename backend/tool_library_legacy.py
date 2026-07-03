import os
import json
import math
import re
import datetime
import requests
import subprocess
import tempfile
import hashlib
from typing import Any, Dict, List, Optional
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from html.parser import HTMLParser
from html import unescape

WORKSPACE_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "workspace")
os.makedirs(WORKSPACE_DIR, exist_ok=True)


def _dedup_results(all_results: List[Dict]) -> List[Dict]:
    seen_urls = set()
    deduped = []
    for r in all_results:
        url = r.get("href", "").rstrip("/").lower()
        if url and url not in seen_urls:
            seen_urls.add(url)
            deduped.append(r)
    return deduped


def _search_duckduckgo(query: str, num_results: int) -> List[Dict]:
    try:
        from ddgs import DDGS
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=num_results, region="us-en"))
        return [{"title": r.get("title", ""), "href": r.get("href", ""), "body": r.get("body", ""), "engine": "duckduckgo"} for r in results]
    except Exception:
        return []


def _search_bing_scrape(query: str, num_results: int) -> List[Dict]:
    try:
        import re as _re
        from html import unescape
        headers = {
            "User-Agent": "Mozilla/5.0 (iPhone; CPU iPhone OS 16_6 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/16.6 Mobile/15E148 Safari/604.1",
            "Accept": "text/html",
            "Accept-Language": "en-US,en;q=0.9",
        }
        resp = requests.get(
            "https://www.bing.com/search",
            params={"q": query, "count": num_results},
            headers=headers,
            timeout=10,
        )
        if resp.status_code != 200:
            return []
        html = resp.text
        results = []
        seen_urls = set()
        links = _re.findall(r'href="(https?://[^"]+)"', html)
        for url in links:
            if "bing.com" in url or "microsoft.com" in url or "go.microsoft" in url:
                continue
            if url in seen_urls:
                continue
            seen_urls.add(url)
            title = url.split("//")[-1].split("/")[0]
            results.append({"title": title, "href": url, "body": "", "engine": "bing"})
            if len(results) >= num_results:
                break
        return results
    except Exception:
        return []


def _search_brave_api(query: str, num_results: int, api_key: str) -> List[Dict]:
    try:
        resp = requests.get(
            "https://api.search.brave.com/res/v1/web/search",
            params={"q": query, "count": num_results},
            headers={"X-Subscription-Token": api_key, "Accept": "application/json"},
            timeout=10,
        )
        if resp.status_code != 200:
            return []
        data = resp.json()
        results = []
        for r in data.get("web", {}).get("results", []):
            results.append({
                "title": r.get("title", ""),
                "href": r.get("url", ""),
                "body": r.get("description", ""),
                "engine": "brave",
            })
        return results
    except Exception:
        return []


def _search_searxng(query: str, num_results: int, base_url: str) -> List[Dict]:
    try:
        resp = requests.get(
            f"{base_url.rstrip('/')}/search",
            params={"q": query, "format": "json", "engines": "google,bing,duckduckgo"},
            timeout=15,
        )
        if resp.status_code != 200:
            return []
        data = resp.json()
        results = []
        for r in data.get("results", []):
            results.append({
                "title": r.get("title", ""),
                "href": r.get("url", ""),
                "body": r.get("content", ""),
                "engine": r.get("engine", "searxng"),
            })
        return results[:num_results]
    except Exception:
        return []


def _optimize_search_query(query: str) -> str:
    """Convert natural language questions into search-engine-friendly keywords."""
    q = query.strip()
    prefixes = [
        r"^how\s+(do|can|to|would)\s+(i|you|we)\s+",
        r"^how\s+",
        r"^what\s+(is|are|do|does|can|should|would)\s+(a|an|the|i|you|we)?\s*",
        r"^what\s+",
        r"^why\s+(do|does|is|are|can|should)\s+(i|you|we)?\s*",
        r"^why\s+",
        r"^when\s+(do|does|is|are|can|should)\s+(i|you|we)?\s*",
        r"^when\s+",
        r"^where\s+(do|does|is|are|can|should)\s+(i|you|we)?\s*",
        r"^where\s+",
        r"^which\s+(is|are|do|does|can|should)\s+(a|an|the|i|you|we)?\s*",
        r"^which\s+",
        r"^can\s+(i|you|we)\s+",
        r"^can\s+",
        r"^should\s+(i|you|we)\s+",
        r"^should\s+",
        r"^is\s+it\s+",
        r"^is\s+",
        r"^are\s+",
        r"^do\s+(i|you|we)\s+",
        r"^does\s+",
    ]
    import re as _re
    for pattern in prefixes:
        q = _re.sub(pattern, "", q, count=1, flags=_re.IGNORECASE)
    stopwords = {"a", "an", "the", "to", "for", "of", "in", "on", "at", "by", "with",
                 "from", "this", "that", "it", "its", "and", "or", "but", "not", "very",
                 "really", "just", "also", "so", "too", "my", "your", "our", "their",
                 "downloaded", "download", "using", "use", "get", "got", "make", "made",
                 "way", "best", "good", "proper", "correct", "right"}
    words = [w for w in q.split() if w.lower() not in stopwords]
    if len(words) > 5:
        words = words[:5]
    return " ".join(words) if words else query


def _is_irrelevant_result(result: Dict, query: str) -> bool:
    """Filter out obviously irrelevant search results."""
    text = (result.get("href", "") + " " + result.get("title", "") + " " + result.get("body", "")).lower()
    irrelevant_domains = ["poki.com", "friv.com", "crazygames.com", "miniclip.com",
                          "kongregate.com", "newgrounds.com", "armor-games.com",
                          "play.google.com/store/apps", "apps.apple.com",
                          "amazon.com", "ebay.com", "walmart.com", "aliexpress.com",
                          "tiktok.com", "instagram.com", "facebook.com",
                          "netflix.com", "hulu.com", "disneyplus.com",
                          "softonic.com", "uptodown.com", "filehippo.com",
                          "trustpilot.com", "glassdoor.com", "reddit.com/r/gaming"]
    for domain in irrelevant_domains:
        if domain in text:
            return True
    irrelevant_phrases = ["play free games", "online games", "free online games",
                          "survival shooter", "battle royale", "mobile game",
                          "play now", "unblocked games", "io games",
                          "free fire", "freegames", "free.fr"]
    score = sum(1 for kw in irrelevant_phrases if kw in text)
    if score >= 1:
        return True
    query_words = set(query.lower().split())
    query_words -= {"how", "do", "you", "to", "the", "a", "an", "is", "are", "what",
                    "why", "when", "where", "which", "can", "should", "for", "of",
                    "in", "on", "at", "and", "or", "it", "this", "that", "i", "we"}
    result_text = result.get("title", "").lower() + " " + result.get("body", "").lower()
    matching_words = sum(1 for w in query_words if w in result_text)
    if query_words and matching_words == 0:
        return True
    return False


class _HTMLTextExtractor(HTMLParser):
    """Strip HTML tags and return clean text. Firecrawl-style content extraction."""

    SKIP_TAGS = {"script", "style", "noscript", "iframe", "svg", "nav", "footer",
                 "header", "aside", "form", "button", "input", "select", "textarea",
                 "meta", "link", "head"}
    BLOCK_TAGS = {"p", "div", "br", "h1", "h2", "h3", "h4", "h5", "h6", "li",
                  "tr", "blockquote", "pre", "section", "article", "main", "figcaption"}
    TEXT_TAGS = {"a": "href", "img": "alt", "source": "src"}

    def __init__(self):
        super().__init__()
        self._text_parts: List[str] = []
        self._skip_depth = 0
        self._skip_tag = ""

    def handle_starttag(self, tag: str, attrs: List[tuple]):
        tag_lower = tag.lower()
        if tag_lower in self.SKIP_TAGS:
            self._skip_depth += 1
            self._skip_tag = tag_lower
            return
        if tag_lower in self.BLOCK_TAGS:
            self._text_parts.append("\n")
        if tag_lower == "a":
            attr_dict = dict(attrs)
            href = attr_dict.get("href", "")
            if href and href.startswith(("http://", "https://")):
                self._text_parts.append(f" [{href}] ")
        if tag_lower == "img":
            attr_dict = dict(attrs)
            alt = attr_dict.get("alt", "")
            if alt:
                self._text_parts.append(f" [{alt}] ")

    def handle_endtag(self, tag: str):
        tag_lower = tag.lower()
        if tag_lower in self.SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1
            return
        if tag_lower in self.BLOCK_TAGS:
            self._text_parts.append("\n")

    def handle_data(self, data: str):
        if self._skip_depth > 0:
            return
        text = data.strip()
        if text:
            self._text_parts.append(text + " ")

    def get_text(self) -> str:
        raw = " ".join(self._text_parts)
        raw = unescape(raw)
        raw = re.sub(r'[ \t]+', ' ', raw)
        raw = re.sub(r'\n{3,}', '\n\n', raw)
        return raw.strip()


def web_scraper(url: str, max_chars=8000, **kwargs) -> str:
    """Firecrawl-clone: fetch a URL, strip all HTML/code, return clean text."""
    try:
        max_chars = int(max_chars) if max_chars else 8000
    except (ValueError, TypeError):
        max_chars = 8000
    try:
        if not url.startswith(("http://", "https://")):
            url = "https://" + url

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                          "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            "Accept-Language": "en-US,en;q=0.9",
            "Accept-Encoding": "gzip, deflate",
        }

        resp = requests.get(url, headers=headers, timeout=15, allow_redirects=True)
        resp.raise_for_status()

        content_type = resp.headers.get("Content-Type", "")
        if "text/html" not in content_type and "application/xhtml" not in content_type:
            if "text/" in content_type:
                return resp.text[:max_chars]
            return f"Non-HTML content ({content_type}). Cannot scrape."

        html = resp.text

        extractor = _HTMLTextExtractor()
        extractor.feed(html)
        clean_text = extractor.get_text()

        title_match = re.search(r'<title[^>]*>(.*?)</title>', html, re.IGNORECASE | re.DOTALL)
        title = title_match.group(1).strip() if title_match else ""
        title = unescape(re.sub(r'<[^>]+>', '', title)).strip()

        if title:
            clean_text = f"# {title}\n\n{clean_text}"

        if len(clean_text) > max_chars:
            clean_text = clean_text[:max_chars] + "\n\n... [content truncated]"

        if not clean_text.strip():
            return f"Scraped {url} but extracted text was empty (likely JS-rendered page)."

        return clean_text

    except requests.exceptions.Timeout:
        return f"Error: Timed out scraping {url}"
    except requests.exceptions.ConnectionError:
        return f"Error: Could not connect to {url}"
    except requests.exceptions.HTTPError as e:
        return f"Error: HTTP {e.response.status_code} from {url}"
    except Exception as e:
        return f"Error scraping {url}: {str(e)}"


def _scrape_search_result_urls(urls: List[str], max_chars_per_page: int = 4000) -> Dict[str, str]:
    """Scrape multiple URLs in parallel and return {url: clean_text}."""
    results: Dict[str, str] = {}
    if not urls:
        return results

    def _fetch_one(url: str) -> tuple:
        text = web_scraper(url, max_chars=max_chars_per_page)
        return (url, text)

    with ThreadPoolExecutor(max_workers=3) as executor:
        futures = {executor.submit(_fetch_one, u): u for u in urls[:3]}
        for future in as_completed(futures, timeout=20):
            try:
                url, text = future.result()
                results[url] = text
            except Exception:
                results[futures[future]] = "Error: Failed to scrape"

    return results


def web_search(query: str, num_results=5, scrape=True) -> str:
    try:
        num_results = int(num_results) if num_results else 5
        scrape = str(scrape).lower() == "true" if isinstance(scrape, str) else bool(scrape)
    except (ValueError, TypeError):
        num_results = 5
        scrape = True
    try:
        from settings_manager import settings_manager
        search_settings = settings_manager.get_search_settings()
    except Exception:
        search_settings = {"search_provider": "auto", "brave_api_key": "", "searxng_url": ""}

    provider = search_settings.get("search_provider", "auto")
    brave_key = search_settings.get("brave_api_key", "")
    searxng_url = search_settings.get("searxng_url", "")

    optimized_query = _optimize_search_query(query)

    fetch_count = max(num_results * 4, 20)
    all_results: List[Dict] = []

    if provider == "brave" and brave_key:
        all_results = _search_brave_api(optimized_query, fetch_count, brave_key)
    elif provider == "searxng" and searxng_url:
        all_results = _search_searxng(optimized_query, fetch_count, searxng_url)
    elif provider == "duckduckgo":
        all_results = _search_duckduckgo(optimized_query, fetch_count)
    else:
        engines = []
        if brave_key:
            engines.append(("brave", lambda: _search_brave_api(optimized_query, fetch_count, brave_key)))
        if searxng_url:
            engines.append(("searxng", lambda: _search_searxng(optimized_query, fetch_count, searxng_url)))
        engines.append(("duckduckgo", lambda: _search_duckduckgo(optimized_query, fetch_count)))
        engines.append(("bing", lambda: _search_bing_scrape(optimized_query, fetch_count)))

        with ThreadPoolExecutor(max_workers=4) as executor:
            futures = {executor.submit(fn): name for name, fn in engines}
            for future in as_completed(futures, timeout=20):
                engine_name = futures[future]
                try:
                    results = future.result()
                    all_results.extend(results)
                except Exception:
                    pass

    all_results = _dedup_results(all_results)
    all_results = [r for r in all_results if not _is_irrelevant_result(r, optimized_query)]

    if not all_results:
        return f"No results found for: {query}"

    final = all_results[:num_results]

    scraped_content: Dict[str, str] = {}
    if scrape:
        urls_to_scrape = [r.get("href", "") for r in final if r.get("href", "").startswith("http")]
        if urls_to_scrape:
            scraped_content = _scrape_search_result_urls(urls_to_scrape, max_chars_per_page=4000)

    output_lines = []
    engines_used = set()
    for i, r in enumerate(final, 1):
        engines_used.add(r.get("engine", "unknown"))
        url = r.get("href", "N/A")
        output_lines.append(f"{i}. {r.get('title', 'No title')}")
        output_lines.append(f"   URL: {url}")
        output_lines.append(f"   {r.get('body', 'No description')}")
        if url in scraped_content and scraped_content[url]:
            page_text = scraped_content[url]
            output_lines.append(f"   [Scraped content]: {page_text[:2000]}")
        output_lines.append("")
    output_lines.append(f"--- Sources: {', '.join(sorted(engines_used))} ---")
    return "\n".join(output_lines)


def read_file(file_path: str, **kwargs) -> str:
    try:
        path = Path(file_path)
        if not path.is_absolute():
            path = Path(WORKSPACE_DIR) / file_path
        if not path.exists():
            return f"Error: File not found: {file_path}"
        suffix = path.suffix.lower()
        if suffix == ".txt" or suffix == ".md" or suffix == ".py" or suffix == ".json" or suffix == ".csv":
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            if len(content) > 10000:
                content = content[:10000] + "\n... [truncated]"
            return content
        elif suffix == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(str(path)) as pdf:
                    text_parts = []
                    for i, page in enumerate(pdf.pages[:30]):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            text_parts.append(f"--- Page {i + 1} ---\n{page_text}")
                    return "\n\n".join(text_parts)[:10000]
            except ImportError:
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(str(path))
                    text_parts = []
                    for page in reader.pages[:20]:
                        text_parts.append(page.extract_text() or "")
                    return "\n".join(text_parts)[:10000]
                except ImportError:
                    return "Error: Neither pdfplumber nor pypdf installed"
        elif suffix == ".docx":
            try:
                from docx import Document
                doc = Document(str(path))
                text_parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                for table in doc.tables:
                    text_parts.append("[Table]")
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        text_parts.append(" | ".join(cells))
                return "\n".join(text_parts)[:10000]
            except ImportError:
                return "Error: python-docx not installed"
        elif suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                text_parts = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = [f"--- Slide {slide_num} ---"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    slide_text.append(para.text)
                        if shape.has_table:
                            slide_text.append("[Table]")
                            for row in shape.table.rows:
                                cells = [cell.text.strip() for cell in row.cells]
                                slide_text.append(" | ".join(cells))
                    if len(slide_text) > 1:
                        text_parts.append("\n".join(slide_text))
                return "\n\n".join(text_parts)[:10000]
            except ImportError:
                return "Error: python-pptx not installed"
        elif suffix in [".xlsx", ".xls"]:
            try:
                import pandas as pd
                df = pd.read_excel(str(path))
                return df.to_string()[:10000]
            except ImportError:
                return "Error: openpyxl not installed"
        else:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()[:10000]
    except Exception as e:
        return f"Error reading file: {str(e)}"


def write_file(file_path: str, content: str, **kwargs) -> str:
    try:
        path = Path(file_path)
        if not path.is_absolute():
            path = Path(WORKSPACE_DIR) / file_path
        path.parent.mkdir(parents=True, exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return f"Successfully wrote {len(content)} bytes to {path}"
    except Exception as e:
        return f"Error writing file: {str(e)}"


DANGEROUS_MODULES = {
    "subprocess", "shutil", "ctypes", "importlib", "sys", "os",
    "socket", "http", "ftplib", "smtplib", "telnetlib",
    "xmlrpc", "pickle", "shelve", "dbm",
}

DANGEROUS_COMMANDS = {
    "rm", "rmdir", "rd", "del", "format", "mkfs",
    "shutdown", "reboot", "halt", "poweroff",
    "sudo", "su", "runas", "net user", "net localgroup",
    "reg delete", "reg add",
    "taskkill", "taskkill /f",
    "cipher", "icacls",
    "bcdboot", "bootrec",
}

SAFE_EVAL_GLOBALS = {"__builtins__": {}}
SAFE_EVAL_LOCALS = {
    "True": True, "False": False, "None": None,
    "len": len, "str": str, "int": int, "float": float,
    "bool": bool, "abs": abs, "min": min, "max": max,
    "round": round, "sum": sum, "sorted": sorted,
}


def safe_eval(expression: str, extra_vars: Optional[Dict] = None) -> Any:
    """Evaluate an expression with restricted builtins and whitelisted names."""
    import ast
    try:
        tree = ast.parse(expression, mode="eval")
    except SyntaxError:
        raise ValueError(f"Invalid expression syntax: {expression}")
    for node in ast.walk(tree):
        if isinstance(node, ast.Call):
            raise ValueError("Function calls are not allowed in expressions")
        if isinstance(node, (ast.Import, ast.ImportFrom)):
            raise ValueError("Imports are not allowed in expressions")
    locals_dict = dict(SAFE_EVAL_LOCALS)
    if extra_vars:
        locals_dict.update(extra_vars)
    return eval(expression, SAFE_EVAL_GLOBALS, locals_dict)


def _check_code_safety(code: str) -> Optional[str]:
    """Check code for dangerous imports/operations. Returns warning or None."""
    code_lower = code.lower()
    for module in DANGEROUS_MODULES:
        patterns = [f"import {module}", f"from {module}", f"__import__('{module}')"]
        for pattern in patterns:
            if pattern in code_lower:
                return f"Warning: code uses restricted module '{module}'"
    return None


def _check_command_safety(command: str) -> Optional[str]:
    """Check command for dangerous operations. Returns warning or None."""
    cmd_lower = command.lower().strip()
    for dangerous in DANGEROUS_COMMANDS:
        if cmd_lower.startswith(dangerous) or f" {dangerous} " in cmd_lower:
            return f"Warning: command contains potentially dangerous operation '{dangerous}'"
    return None


def run_code(code: str, **kwargs) -> str:
    warning = _check_code_safety(code)
    try:
        with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False, dir=WORKSPACE_DIR) as f:
            f.write(code)
            temp_path = f.name
        result = subprocess.run(
            ["python", temp_path],
            capture_output=True,
            text=True,
            timeout=30,
            cwd=WORKSPACE_DIR
        )
        os.unlink(temp_path)
        output = result.stdout
        if result.stderr:
            output += "\nSTDERR:\n" + result.stderr
        if not output.strip():
            output = "Code executed successfully (no output)."
        if warning:
            output = f"[SAFETY] {warning}\n\n{output}"
        return output[:5000]
    except subprocess.TimeoutExpired:
        return "Error: Code execution timed out (30s limit)"
    except Exception as e:
        return f"Error running code: {str(e)}"


def send_email(to: str, subject: str, body: str, smtp_server: str = "smtp.gmail.com",
               smtp_port: int = 587, username: str = "", password: str = "", **kwargs) -> str:
    try:
        import smtplib
        from email.mime.text import MIMEText
        msg = MIMEText(body)
        msg["Subject"] = subject
        msg["From"] = username
        msg["To"] = to
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()
            if username and password:
                server.login(username, password)
            server.send_message(msg)
        return f"Email sent to {to}"
    except Exception as e:
        return f"Error sending email: {str(e)}"


def http_request(url: str, method: str = "GET", headers: str = "",
                 body: str = "", timeout: int = 30, **kwargs) -> str:
    try:
        parsed_headers = {}
        if headers:
            try:
                parsed_headers = json.loads(headers) if isinstance(headers, str) else headers
            except json.JSONDecodeError:
                for line in headers.split("\n"):
                    if ":" in line:
                        k, v = line.split(":", 1)
                        parsed_headers[k.strip()] = v.strip()
        kwargs_req = {
            "url": url,
            "headers": parsed_headers,
            "timeout": timeout
        }
        if method.upper() == "POST":
            kwargs_req["data"] = body
        elif method.upper() == "PUT":
            kwargs_req["data"] = body
        response = requests.request(method.upper(), **kwargs_req)
        content = response.text[:10000]
        return f"Status: {response.status_code}\n\n{content}"
    except Exception as e:
        return f"HTTP request error: {str(e)}"


def calculate(expression: str, **kwargs) -> str:
    try:
        result = safe_eval(expression, {
            "pi": math.pi, "e": math.e, "sqrt": math.sqrt,
            "pow": pow, "log": math.log, "sin": math.sin,
            "cos": math.cos, "tan": math.tan, "ceil": math.ceil,
            "floor": math.floor
        })
        return str(result)
    except Exception as e:
        return f"Calculation error: {str(e)}"


def run_command(command: str, working_directory: str = "", timeout: int = 60, **kwargs) -> str:
    warning = _check_command_safety(command)
    timeout = min(timeout, 120)
    try:
        import sys
        is_windows = sys.platform == "win32"
        if is_windows:
            shell_cmd = ["cmd", "/c", command]
        else:
            shell_cmd = ["bash", "-c", command]
        if working_directory and Path(working_directory).is_dir():
            cwd = working_directory
        else:
            cwd = WORKSPACE_DIR
        result = subprocess.run(
            shell_cmd,
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=cwd
        )
        output = ""
        if result.stdout:
            output += result.stdout
        if result.stderr:
            output += ("\n" if output else "") + result.stderr
        if not output.strip():
            output = "Command executed successfully (no output)."
        exit_code = result.returncode
        if warning:
            output = f"[SAFETY] {warning}\n\n{output}"
        return f"Exit Code: {exit_code}\n\n{output[:8000]}"
    except subprocess.TimeoutExpired:
        return f"Error: Command timed out after {timeout}s"
    except FileNotFoundError:
        return f"Error: Shell not found. Command: {command}"
    except Exception as e:
        return f"Error running command: {str(e)}"


def get_datetime(format_str: str = "%Y-%m-%d %H:%M:%S", **kwargs) -> str:
    return datetime.datetime.now().strftime(format_str)


def _get_browser_path() -> str:
    """Get browser executable path from settings or env var."""
    try:
        from settings_manager import settings_manager
        path = settings_manager.get("browser_path", "")
        if path and os.path.exists(path):
            return path
    except Exception:
        pass
    env_path = os.environ.get("OLLAMAFLOW_BROWSER_PATH", "")
    if env_path and os.path.exists(env_path):
        return env_path
    brave_paths = [
        r"C:\Program Files\BraveSoftware\Brave-Browser\Application\brave.exe",
        r"C:\Program Files (x86)\BraveSoftware\Brave-Browser\Application\brave.exe",
        os.path.expanduser(r"~\AppData\Local\BraveSoftware\Brave-Browser\Application\brave.exe"),
        "/usr/bin/brave-browser",
        "/usr/bin/brave",
        "/Applications/Brave Browser.app/Contents/MacOS/Brave Browser",
    ]
    for p in brave_paths:
        if os.path.exists(p):
            return p
    return ""


def playwright_browser(action: str = "goto", url: str = "", browser: str = "chrome",
                       selector: str = "", text: str = "", screenshot: str = "",
                       wait_seconds: int = 3, **kwargs) -> str:
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser_type = p.chromium
            if browser == "brave":
                brave_path = _get_browser_path()
                if not brave_path:
                    return "Error: Brave browser not found. Set browser_path in Settings or OLLAMAFLOW_BROWSER_PATH env var."
                launch_opts = {"headless": True, "executable_path": brave_path}
            elif browser == "msedge":
                launch_opts = {"headless": True, "channel": "msedge"}
            elif browser == "firefox":
                browser_type = p.firefox
                launch_opts = {"headless": True}
            elif browser == "webkit":
                browser_type = p.webkit
                launch_opts = {"headless": True}
            else:
                launch_opts = {"headless": True, "channel": "chrome"}

            pw_browser = browser_type.launch(**launch_opts)
            page = pw_browser.new_page()

            if url and not url.startswith(("http://", "https://")):
                url = "https://" + url

            if action == "goto" and url:
                page.goto(url, timeout=30000)
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(2000)
                title = page.title()
                content = page.content()[:15000]
                pw_browser.close()
                return f"Title: {title}\n\n{content}"

            elif action == "click" and selector:
                page.goto(url, timeout=30000) if url else None
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(2000)
                page.click(selector)
                page.wait_for_timeout(wait_seconds * 1000)
                content = page.content()[:15000]
                pw_browser.close()
                return f"Clicked: {selector}\n\n{content}"

            elif action == "type" and selector and text:
                page.goto(url, timeout=30000) if url else None
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(2000)
                page.fill(selector, text)
                content = page.content()[:15000]
                pw_browser.close()
                return f"Typed into: {selector}\n\n{content}"

            elif action == "extract" and selector:
                page.goto(url, timeout=30000) if url else None
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(2000)
                elements = page.query_selector_all(selector)
                texts = [el.inner_text() for el in elements[:20]]
                pw_browser.close()
                return f"Extracted {len(texts)} elements:\n" + "\n---\n".join(texts)

            elif action == "screenshot":
                page.goto(url, timeout=30000) if url else None
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(3000)
                screenshot_path = screenshot or "screenshot.png"
                if not os.path.isabs(screenshot_path):
                    screenshot_path = os.path.join(WORKSPACE_DIR, screenshot_path)
                page.screenshot(path=screenshot_path, full_page=True)
                pw_browser.close()
                return f"Screenshot saved to: {screenshot_path}"

            elif action == "evaluate" and text:
                page.goto(url, timeout=30000) if url else None
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(2000)
                result = page.evaluate(text)
                pw_browser.close()
                return f"Result: {str(result)[:5000]}"

            else:
                pw_browser.close()
                return f"Error: Invalid action '{action}' or missing parameters. Actions: goto, click, type, extract, screenshot, evaluate"

    except Exception as e:
        return f"Browser error: {str(e)}"


def list_directory(path: str = "", pattern: str = "", **kwargs) -> str:
    try:
        dir_path = Path(path) if path else Path(WORKSPACE_DIR)
        if not dir_path.is_absolute():
            dir_path = Path(WORKSPACE_DIR) / path
        if not dir_path.exists():
            return f"Error: Directory not found: {path}"
        if not dir_path.is_dir():
            return f"Error: Not a directory: {path}"
        entries = []
        glob_pattern = pattern if pattern else "*"
        for item in sorted(dir_path.glob(glob_pattern)):
            if item.name.startswith("."):
                continue
            try:
                size = item.stat().st_size
                if size < 1024:
                    size_str = f"{size}B"
                elif size < 1024 * 1024:
                    size_str = f"{size / 1024:.1f}KB"
                else:
                    size_str = f"{size / (1024 * 1024):.1f}MB"
            except Exception:
                size_str = "?"
            item_type = "DIR " if item.is_dir() else "FILE"
            entries.append(f"  [{item_type}] {item.name:40s} {size_str}")
        if not entries:
            return f"Directory is empty: {dir_path}"
        header = f"Contents of {dir_path}:\n"
        return header + "\n".join(entries)
    except Exception as e:
        return f"Error listing directory: {str(e)}"


def search_files(query: str, path: str = "", search_type: str = "name", **kwargs) -> str:
    try:
        search_dir = Path(path) if path else Path(WORKSPACE_DIR)
        if not search_dir.is_absolute():
            search_dir = Path(WORKSPACE_DIR) / path
        if not search_dir.exists():
            return f"Error: Directory not found: {path}"
        results = []
        if search_type == "name":
            for item in search_dir.rglob(f"*{query}*"):
                if item.name.startswith(".") or "node_modules" in str(item):
                    continue
                rel = item.relative_to(search_dir)
                results.append(str(rel))
                if len(results) >= 50:
                    break
        elif search_type == "content":
            pattern = re.compile(query, re.IGNORECASE)
            for item in search_dir.rglob("*"):
                if item.name.startswith(".") or "node_modules" in str(item):
                    continue
                if not item.is_file():
                    continue
                try:
                    text = item.read_text(encoding="utf-8", errors="replace")[:50000]
                    for i, line in enumerate(text.split("\n"), 1):
                        if pattern.search(line):
                            rel = item.relative_to(search_dir)
                            results.append(f"{rel}:{i}: {line.strip()[:120]}")
                            if len(results) >= 50:
                                break
                except Exception:
                    continue
                if len(results) >= 50:
                    break
        if not results:
            return f"No results found for '{query}' in {search_dir}"
        return f"Found {len(results)} results:\n" + "\n".join(results)
    except Exception as e:
        return f"Error searching files: {str(e)}"


def json_query(data: str, query: str, **kwargs) -> str:
    try:
        if isinstance(data, str):
            parsed = json.loads(data)
        else:
            parsed = data
        parts = [p for p in re.split(r'\.(?![^\[]*\])', query) if p]
        current = parsed
        for part in parts:
            array_match = re.match(r'^(\w+)\[(\d+)\]$', part)
            if array_match:
                key = array_match.group(1)
                idx = int(array_match.group(2))
                if isinstance(current, dict):
                    current = current[key]
                current = current[idx]
            elif isinstance(current, dict):
                current = current[part]
            elif isinstance(current, list):
                current = current[int(part)]
            else:
                return f"Error: Cannot navigate into {type(current).__name__} with key '{part}'"
        if isinstance(current, (dict, list)):
            return json.dumps(current, indent=2, ensure_ascii=False)[:10000]
        return str(current)
    except (KeyError, IndexError, TypeError) as e:
        return f"Error: Query '{query}' failed - {str(e)}"
    except json.JSONDecodeError:
        return "Error: Invalid JSON input"
    except Exception as e:
        return f"Error: {str(e)}"


def csv_analyze(file_path: str, operation: str = "summary", column: str = "",
                filter_value: str = "", limit: int = 20, **kwargs) -> str:
    try:
        import pandas as pd
        path = Path(file_path)
        if not path.is_absolute():
            path = Path(WORKSPACE_DIR) / file_path
        if not path.exists():
            return f"Error: File not found: {file_path}"
        df = pd.read_csv(str(path))
        if operation == "summary":
            buf = []
            buf.append(f"Shape: {df.shape[0]} rows x {df.shape[1]} columns")
            buf.append(f"\nColumns: {', '.join(df.columns.tolist())}")
            buf.append(f"\nData types:\n{df.dtypes.to_string()}")
            buf.append(f"\nBasic stats:\n{df.describe().to_string()}")
            buf.append(f"\nNull counts:\n{df.isnull().sum().to_string()}")
            return "\n".join(buf)[:10000]
        elif operation == "head":
            return df.head(limit).to_string()[:10000]
        elif operation == "filter" and column and filter_value:
            try:
                filtered = df[df[column].astype(str).str.contains(filter_value, case=False, na=False)]
            except Exception:
                filtered = df[df[column] == filter_value]
            return f"Filtered ({len(filtered)} rows):\n{filtered.head(limit).to_string()}"[:10000]
        elif operation == "sort" and column:
            ascending = not filter_value.lower().startswith("desc")
            sorted_df = df.sort_values(by=column, ascending=ascending)
            return sorted_df.head(limit).to_string()[:10000]
        elif operation == "group" and column:
            grouped = df.groupby(column).size().reset_index(name="count")
            grouped = grouped.sort_values("count", ascending=False)
            return grouped.head(limit).to_string()[:10000]
        else:
            return f"Error: Invalid operation '{operation}' or missing column parameter"
    except ImportError:
        return "Error: pandas not installed"
    except Exception as e:
        return f"Error analyzing CSV: {str(e)}"


def database_query(db_path: str, query: str, params: str = "", **kwargs) -> str:
    try:
        import sqlite3
        path = Path(db_path)
        if not path.is_absolute():
            path = Path(WORKSPACE_DIR) / db_path
        if not path.exists():
            return f"Error: Database not found: {db_path}"
        bind_params = []
        if params:
            try:
                bind_params = json.loads(params) if isinstance(params, str) else params
            except Exception:
                pass
        conn = sqlite3.connect(str(path))
        try:
            cursor = conn.cursor()
            if bind_params:
                cursor.execute(query, bind_params)
            else:
                cursor.execute(query)
            if query.strip().upper().startswith(("INSERT", "UPDATE", "DELETE", "DROP", "CREATE", "ALTER")):
                conn.commit()
                return f"Query executed. Rows affected: {cursor.rowcount}"
            rows = cursor.fetchall()
            if not rows:
                return "Query returned no results."
            col_names = [desc[0] for desc in cursor.description] if cursor.description else []
            lines = [" | ".join(str(c) for c in col_names)]
            lines.append("-" * len(lines[0]))
            for row in rows[:100]:
                lines.append(" | ".join(str(v) for v in row))
            result = "\n".join(lines)
            if len(rows) > 100:
                result += f"\n... ({len(rows) - 100} more rows)"
            return result[:10000]
        finally:
            conn.close()
    except Exception as e:
        return f"Database error: {str(e)}"


def random_generate(type: str = "uuid", length: int = 16, count: int = 1, **kwargs) -> str:
    try:
        import uuid as uuid_mod
        import secrets
        import string
        results = []
        for _ in range(min(count, 100)):
            if type == "uuid":
                results.append(str(uuid_mod.uuid4()))
            elif type == "password":
                alphabet = string.ascii_letters + string.digits + "!@#$%^&*"
                results.append("".join(secrets.choice(alphabet) for _ in range(max(length, 8))))
            elif type == "number":
                upper = 10 ** min(length, 15)
                results.append(str(secrets.randbelow(upper)))
            elif type == "string":
                results.append("".join(secrets.choice(string.ascii_lowercase + string.digits) for _ in range(max(length, 1))))
            else:
                return f"Error: Unknown type '{type}'. Use: uuid, password, number, string"
        return "\n".join(results)
    except Exception as e:
        return f"Error generating random data: {str(e)}"


def diff_texts(text_a: str, text_b: str, context_lines: int = 3, **kwargs) -> str:
    try:
        import difflib
        lines_a = text_a.splitlines(keepends=True)
        lines_b = text_b.splitlines(keepends=True)
        diff = list(difflib.unified_diff(
            lines_a, lines_b,
            fromfile="text_a", tofile="text_b",
            n=context_lines
        ))
        if not diff:
            return "No differences found."
        return "".join(diff)[:10000]
    except Exception as e:
        return f"Error computing diff: {str(e)}"


def hash_data(data: str, algorithm: str = "sha256", **kwargs) -> str:
    try:
        alg_map = {
            "sha256": hashlib.sha256,
            "sha1": hashlib.sha1,
            "md5": hashlib.md5,
        }
        if algorithm not in alg_map:
            return f"Error: Unknown algorithm '{algorithm}'. Use: sha256, sha1, md5"
        h = alg_map[algorithm]()
        h.update(data.encode("utf-8"))
        return f"{algorithm}: {h.hexdigest()}"
    except Exception as e:
        return f"Error hashing data: {str(e)}"


def rss_read(url: str, max_items: int = 10, **kwargs) -> str:
    try:
        import feedparser
        feed = feedparser.parse(url)
        if feed.bozo and not feed.entries:
            return f"Error parsing feed: {feed.bozo_exception}"
        title = feed.feed.get("title", "Unknown Feed")
        lines = [f"Feed: {title}", f"URL: {url}", ""]
        for i, entry in enumerate(feed.entries[:max_items], 1):
            entry_title = entry.get("title", "No title")
            link = entry.get("link", "")
            published = entry.get("published", entry.get("updated", ""))
            summary = entry.get("summary", "")[:200]
            lines.append(f"{i}. {entry_title}")
            if published:
                lines.append(f"   Date: {published}")
            if link:
                lines.append(f"   Link: {link}")
            if summary:
                lines.append(f"   {summary}")
            lines.append("")
        return "\n".join(lines)[:10000]
    except ImportError:
        return "Error: feedparser not installed. Run: pip install feedparser"
    except Exception as e:
        return f"Error reading RSS feed: {str(e)}"


def slack_webhook(webhook_url: str, message: str, channel: str = "",
                  username: str = "OllamaFlow", **kwargs) -> str:
    try:
        payload = {"text": message, "username": username}
        if channel:
            payload["channel"] = channel
        resp = requests.post(
            webhook_url,
            json=payload,
            timeout=10
        )
        if resp.status_code == 200:
            return f"Message sent successfully via Slack webhook"
        else:
            return f"Slack webhook error: Status {resp.status_code} - {resp.text[:500]}"
    except Exception as e:
        return f"Slack webhook error: {str(e)}"


def web_research(action: str = "search", query: str = "", url: str = "",
                 engine: str = "google", max_results: int = 5, selector: str = "",
                 screenshot: str = "", wait_seconds: int = 3, **kwargs) -> str:
    """Search the web or scrape pages using Playwright with Brave browser to bypass anti-bot defenses."""
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        return "Error: playwright not installed. Run: pip install playwright && playwright install chromium"

    brave_path = _get_browser_path()
    if not brave_path:
        return "Error: Brave browser not found. Set browser_path in Settings or OLLAMAFLOW_BROWSER_PATH env var."

    search_urls = {
        "google": "https://www.google.com/search?q={query}&num={num}",
        "bing": "https://www.bing.com/search?q={query}&count={num}",
        "duckduckgo": "https://duckduckgo.com/?q={query}",
        "brave": "https://search.brave.com/search?q={query}",
    }

    def _extract_page_text(page) -> str:
        return page.evaluate("""() => {
            const skipTags = new Set(['script','style','noscript','iframe','svg','nav','footer','header','aside','form','button','input','select','textarea','meta','link','head']);
            function walk(node) {
                if (node.nodeType === 3) return node.textContent.trim() + ' ';
                if (node.nodeType !== 1) return '';
                if (skipTags.has(node.tagName.toLowerCase())) return '';
                let text = '';
                for (const child of node.childNodes) text += walk(child);
                const tag = node.tagName.toLowerCase();
                if (['p','div','br','h1','h2','h3','h4','h5','h6','li','tr','blockquote','pre','section','article','main'].includes(tag)) text = '\\n' + text + '\\n';
                if (tag === 'a') {
                    const href = node.getAttribute('href');
                    if (href && href.startsWith('http')) text = ' [' + href + '] ' + text;
                }
                return text;
            }
            return walk(document.body).replace(/[ \\t]+/g, ' ').replace(/\\n{3,}/g, '\\n\\n').trim();
        }""")

    def _extract_search_results(page) -> list:
        return page.evaluate("""() => {
            const results = [];
            const seen = new Set();
            const links = document.querySelectorAll('a[href]');
            for (const link of links) {
                const href = link.getAttribute('href');
                if (!href || !href.startsWith('http')) continue;
                if (seen.has(href)) continue;
                const text = link.innerText.trim();
                if (!text || text.length < 5) continue;
                if (href.includes('google.com') || href.includes('bing.com') || href.includes('duckduckgo.com') || href.includes('brave.com')) continue;
                seen.add(href);
                const parent = link.closest('div, li, article');
                let body = '';
                if (parent) {
                    const clone = parent.cloneNode(true);
                    clone.querySelectorAll('a').forEach(a => a.remove());
                    body = clone.innerText.trim().substring(0, 300);
                }
                results.push({title: text, url: href, body: body});
                if (results.length >= 20) break;
            }
            return results;
        }""")

    try:
        with sync_playwright() as p:
            launch_opts = {
                "headless": True,
                "executable_path": brave_path,
                "args": [
                    "--disable-blink-features=AutomationControlled",
                    "--no-first-run",
                    "--no-default-browser-check",
                    "--disable-infobars",
                ]
            }
            browser = p.chromium.launch(**launch_opts)
            context = browser.new_context(
                user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                viewport={"width": 1920, "height": 1080},
                locale="en-US",
            )
            page = context.new_page()
            page.add_init_script("""
                Object.defineProperty(navigator, 'webdriver', {get: () => undefined});
                window.chrome = {runtime: {}};
            """)

            if action == "search":
                if not query:
                    browser.close()
                    return "Error: query parameter is required for search"
                num = max(min(max_results, 20), 1)
                search_url = search_urls.get(engine, search_urls["google"]).format(query=query.replace(" ", "+"), num=num)
                page.goto(search_url, timeout=30000)
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(3000)
                results = _extract_search_results(page)
                browser.close()
                if not results:
                    return f"No search results found for: {query}"
                lines = [f"Search results for: {query} (via {engine})", ""]
                for i, r in enumerate(results[:max_results], 1):
                    lines.append(f"{i}. {r['title']}")
                    lines.append(f"   URL: {r['url']}")
                    if r['body']:
                        lines.append(f"   {r['body']}")
                    lines.append("")
                return "\n".join(lines)

            elif action == "search_and_scrape":
                if not query:
                    browser.close()
                    return "Error: query parameter is required for search_and_scrape"
                num = max(min(max_results, 10), 1)
                search_url = search_urls.get(engine, search_urls["google"]).format(query=query.replace(" ", "+"), num=num)
                page.goto(search_url, timeout=30000)
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(3000)
                results = _extract_search_results(page)
                if not results:
                    browser.close()
                    return f"No search results found for: {query}"
                lines = [f"Search results for: {query} (via {engine})", ""]
                for i, r in enumerate(results[:num], 1):
                    lines.append(f"{i}. {r['title']}")
                    lines.append(f"   URL: {r['url']}")
                    if r['body']:
                        lines.append(f"   {r['body']}")
                    try:
                        page.goto(r['url'], timeout=20000)
                        page.wait_for_load_state("domcontentloaded", timeout=10000)
                        page.wait_for_timeout(2000)
                        text = _extract_page_text(page)
                        if text and len(text) > 50:
                            lines.append(f"   [Scraped]: {text[:3000]}")
                    except Exception:
                        pass
                    lines.append("")
                browser.close()
                return "\n".join(lines)

            elif action == "goto":
                if not url:
                    browser.close()
                    return "Error: url parameter is required for goto"
                if not url.startswith(("http://", "https://")):
                    url = "https://" + url
                page.goto(url, timeout=30000)
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(wait_seconds * 1000)
                title = page.title()
                text = _extract_page_text(page)
                browser.close()
                return f"Title: {title}\n\n{text[:10000]}"

            elif action == "extract":
                if not url:
                    browser.close()
                    return "Error: url parameter is required for extract"
                if not url.startswith(("http://", "https://")):
                    url = "https://" + url
                page.goto(url, timeout=30000)
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(wait_seconds * 1000)
                if selector:
                    elements = page.query_selector_all(selector)
                    texts = [el.inner_text() for el in elements[:20]]
                    browser.close()
                    return f"Extracted {len(texts)} elements matching '{selector}':\n" + "\n---\n".join(texts)
                else:
                    text = _extract_page_text(page)
                    browser.close()
                    return f"Page content:\n\n{text[:10000]}"

            elif action == "screenshot":
                if not url:
                    browser.close()
                    return "Error: url parameter is required for screenshot"
                if not url.startswith(("http://", "https://")):
                    url = "https://" + url
                page.goto(url, timeout=30000)
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(wait_seconds * 1000)
                screenshot_path = screenshot or "screenshot.png"
                if not os.path.isabs(screenshot_path):
                    screenshot_path = os.path.join(WORKSPACE_DIR, screenshot_path)
                page.screenshot(path=screenshot_path, full_page=True)
                browser.close()
                return f"Screenshot saved to: {screenshot_path}"

            elif action == "click":
                if not url or not selector:
                    browser.close()
                    return "Error: url and selector parameters are required for click"
                if not url.startswith(("http://", "https://")):
                    url = "https://" + url
                page.goto(url, timeout=30000)
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(2000)
                page.click(selector)
                page.wait_for_timeout(wait_seconds * 1000)
                text = _extract_page_text(page)
                browser.close()
                return f"Clicked: {selector}\n\n{text[:10000]}"

            elif action == "type":
                if not url or not selector:
                    browser.close()
                    return "Error: url and selector parameters are required for type"
                if not url.startswith(("http://", "https://")):
                    url = "https://" + url
                page.goto(url, timeout=30000)
                page.wait_for_load_state("domcontentloaded", timeout=15000)
                page.wait_for_timeout(2000)
                page.fill(selector, query)
                page.keyboard.press("Enter")
                page.wait_for_timeout(wait_seconds * 1000)
                text = _extract_page_text(page)
                browser.close()
                return f"Typed and submitted in '{selector}'\n\n{text[:10000]}"

            elif action == "evaluate":
                if not query:
                    browser.close()
                    return "Error: query parameter (JS code) is required for evaluate"
                if url:
                    if not url.startswith(("http://", "https://")):
                        url = "https://" + url
                    page.goto(url, timeout=30000)
                    page.wait_for_load_state("domcontentloaded", timeout=15000)
                    page.wait_for_timeout(2000)
                result = page.evaluate(query)
                browser.close()
                return f"Result: {str(result)[:10000]}"

            else:
                browser.close()
                return f"Error: Unknown action '{action}'. Actions: search, search_and_scrape, goto, extract, screenshot, click, type, evaluate"

    except Exception as e:
        return f"Web research error: {str(e)}"


def youtube_transcript(video_url: str, language: str = "en", **kwargs) -> str:
    """Extract transcript from a YouTube video URL."""
    try:
        import re
        video_id_match = re.search(r'(?:v=|youtu\.be/|/embed/)([a-zA-Z0-9_-]{11})', video_url)
        if not video_id_match:
            return f"Error: Could not extract video ID from URL: {video_url}"
        video_id = video_id_match.group(1)
        try:
            from youtube_transcript_api import YouTubeTranscriptApi
            transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=[language])
            lines = [f"[{t['start']:.1f}s] {t['text']}" for t in transcript]
            return "\n".join(lines)[:10000]
        except ImportError:
            pass
        api_url = f"https://www.youtube.com/watch?v={video_id}"
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        resp = requests.get(api_url, headers=headers, timeout=15)
        caption_match = re.search(r'"captions":\s*(\{.*?\})\s*,\s*"videoDetails"', resp.text)
        if not caption_match:
            return f"Error: No captions available for video {video_id}. The video may not have subtitles."
        caption_data = json.loads(caption_match.group(1))
        tracks = caption_data.get("playerCaptionsTracklistRenderer", {}).get("captionTracks", [])
        if not tracks:
            return f"Error: No caption tracks found for video {video_id}."
        track_url = tracks[0].get("baseUrl", "")
        if not track_url:
            return "Error: No caption URL found."
        caption_resp = requests.get(track_url, timeout=10)
        text_matches = re.findall(r'<text[^>]*>(.*?)</text>', caption_resp.text)
        import html as html_mod
        lines = [html_mod.unescape(t) for t in text_matches if t.strip()]
        if not lines:
            return f"Error: Could not parse captions for video {video_id}."
        return "\n".join(lines)[:10000]
    except Exception as e:
        return f"YouTube transcript error: {str(e)}"


def clipboard_copy(text: str = "", **kwargs) -> str:
    """Copy text to the system clipboard."""
    try:
        import subprocess
        import sys
        if sys.platform == "win32":
            process = subprocess.Popen(["clip"], stdin=subprocess.PIPE)
            process.communicate(text.encode("utf-16le"))
        elif sys.platform == "darwin":
            process = subprocess.Popen(["pbcopy"], stdin=subprocess.PIPE)
            process.communicate(text.encode("utf-8"))
        else:
            process = subprocess.Popen(["xclip", "-selection", "clipboard"], stdin=subprocess.PIPE)
            process.communicate(text.encode("utf-8"))
        return f"Copied {len(text)} characters to clipboard."
    except Exception as e:
        return f"Clipboard copy error: {str(e)}"


def clipboard_paste(**kwargs) -> str:
    """Read text from the system clipboard."""
    try:
        import subprocess
        import sys
        if sys.platform == "win32":
            result = subprocess.run(["powershell", "-command", "Get-Clipboard"],
                                    capture_output=True, text=True, timeout=5)
            return result.stdout.strip()
        elif sys.platform == "darwin":
            result = subprocess.run(["pbpaste"], capture_output=True, text=True, timeout=5)
            return result.stdout.strip()
        else:
            result = subprocess.run(["xclip", "-selection", "clipboard", "-o"],
                                    capture_output=True, text=True, timeout=5)
            return result.stdout.strip()
    except Exception as e:
        return f"Clipboard paste error: {str(e)}"


_rate_limit_buckets = {}


def rate_limiter(max_requests: int = 10, window_seconds: int = 60, bucket: str = "default", **kwargs) -> str:
    """Simple in-memory rate limiter. Returns 'allowed' or 'rate_limited'."""
    import time as _time
    now = _time.time()
    if bucket not in _rate_limit_buckets:
        _rate_limit_buckets[bucket] = []
    timestamps = _rate_limit_buckets[bucket]
    timestamps = [t for t in timestamps if now - t < window_seconds]
    if len(timestamps) >= max_requests:
        _rate_limit_buckets[bucket] = timestamps
        wait_time = window_seconds - (now - timestamps[0])
        return f"rate_limited: Wait {wait_time:.1f}s. Limit: {max_requests} requests per {window_seconds}s."
    timestamps.append(now)
    _rate_limit_buckets[bucket] = timestamps
    remaining = max_requests - len(timestamps)
    return f"allowed: {remaining} requests remaining in {window_seconds}s window."


def file_watcher(path: str = "", pattern: str = "*", action: str = "list_changes", since_minutes: int = 60, **kwargs) -> str:
    """Watch or list recent file changes in a directory."""
    try:
        watch_dir = Path(path) if path else Path(WORKSPACE_DIR)
        if not watch_dir.is_absolute():
            watch_dir = Path(WORKSPACE_DIR) / path
        if not watch_dir.exists():
            return f"Error: Directory not found: {path}"
        import time as _time
        cutoff = _time.time() - (since_minutes * 60)
        changes = []
        for item in watch_dir.rglob(pattern):
            if item.name.startswith(".") or "node_modules" in str(item):
                continue
            try:
                mtime = item.stat().st_mtime
                if mtime >= cutoff:
                    age_s = _time.time() - mtime
                    if age_s < 60:
                        age = f"{age_s:.0f}s ago"
                    elif age_s < 3600:
                        age = f"{age_s/60:.0f}m ago"
                    else:
                        age = f"{age_s/3600:.0f}h ago"
                    changes.append({
                        "path": str(item.relative_to(watch_dir)),
                        "type": "dir" if item.is_dir() else "file",
                        "modified": age,
                        "size": item.stat().st_size
                    })
            except Exception:
                continue
        if not changes:
            return f"No changes in last {since_minutes} minutes in {watch_dir}"
        lines = [f"Changes in {watch_dir} (last {since_minutes} min):", ""]
        for c in sorted(changes, key=lambda x: x["path"]):
            size = c["size"]
            if size < 1024:
                size_str = f"{size}B"
            elif size < 1024 * 1024:
                size_str = f"{size/1024:.1f}KB"
            else:
                size_str = f"{size/(1024*1024):.1f}MB"
            lines.append(f"  [{c['type'].upper()}] {c['path']:40s} {size_str:10s} {c['modified']}")
        return "\n".join(lines)[:5000]
    except Exception as e:
        return f"File watcher error: {str(e)}"


def browser_use_tool(task: str, llm_provider: str = "ollama", llm_model: str = "",
                     headless: bool = True, allowed_domains: str = "", max_steps: int = 25,
                     **kwargs) -> str:
    try:
        import asyncio
        from browser_use import Agent, BrowserProfile
        if llm_provider == "ollama":
            from langchain_ollama import ChatOllama
            llm = ChatOllama(model=llm_model or "llama3.1:8b")
        elif llm_provider == "openai":
            from langchain_openai import ChatOpenAI
            llm = ChatOpenAI(model=llm_model or "gpt-4o")
        elif llm_provider == "anthropic":
            from langchain_anthropic import ChatAnthropic
            llm = ChatAnthropic(model=llm_model or "claude-sonnet-4-20250514")
        else:
            return f"Error: Unknown LLM provider '{llm_provider}'. Use ollama, openai, or anthropic."

        domains = [d.strip() for d in allowed_domains.split(",") if d.strip()] if allowed_domains else None
        max_steps = int(max_steps) if max_steps else 25

        async def _run():
            agent = Agent(
                task=task,
                llm=llm,
                browser_profile=BrowserProfile(
                    headless=headless,
                    allowed_domains=domains,
                ),
            )
            history = await agent.run(max_steps=max_steps)
            return history.final_result()

        return asyncio.run(_run())
    except ImportError as e:
        return f"Error: browser-use not installed. Run: pip install browser-use. Details: {e}"
    except Exception as e:
        return f"Browser-use error: {str(e)}"


def crawl4ai_tool(action: str = "scrape", url: str = "", max_pages: int = 5,
                  css_selector: str = "", javascript: str = "", cache: bool = True,
                  fit_markdown: bool = True, **kwargs) -> str:
    try:
        import asyncio
        from crawl4ai import AsyncWebCrawler, BrowserConfig, CrawlerRunConfig, CacheMode
        max_pages = int(max_pages) if max_pages else 5

        async def _crawl():
            browser_config = BrowserConfig(headless=True)
            run_config = CrawlerRunConfig(
                cache_mode=CacheMode.ENABLED if cache else CacheMode.BYPASS,
                css_selector=css_selector or None,
                js_code=javascript or None,
            )

            async with AsyncWebCrawler(config=browser_config) as crawler:
                if action == "scrape":
                    result = await crawler.arun(url=url, config=run_config)
                    if fit_markdown and result.markdown and result.markdown.fit_markdown:
                        return result.markdown.fit_markdown
                    return result.markdown.raw_markdown if result.markdown else "Error: No content returned"

                elif action == "deep_crawl":
                    from crawl4ai.deep_crawling import BFSDeepCrawlStrategy
                    strategy = BFSDeepCrawlStrategy(max_depth=max_pages)
                    run_config.strategy = strategy
                    results = await crawler.arun(url=url, config=run_config)
                    output_parts = []
                    for r in results:
                        md = r.markdown.fit_markdown if (fit_markdown and r.markdown and r.markdown.fit_markdown) else (r.markdown.raw_markdown if r.markdown else "")
                        if md:
                            output_parts.append(f"### {r.url}\n\n{md}")
                    return "\n\n---\n\n".join(output_parts) if output_parts else "Error: No content found during deep crawl"

                elif action == "extract_structured":
                    from crawl4ai import JsonCssExtractionStrategy
                    schema = {"type": "object", "properties": {}, "selector": css_selector or "body"}
                    run_config = CrawlerRunConfig(
                        extraction_strategy=JsonCssExtractionStrategy(schema=schema),
                        cache_mode=CacheMode.ENABLED if cache else CacheMode.BYPASS,
                    )
                    async with AsyncWebCrawler(config=browser_config) as crawler2:
                        result = await crawler2.arun(url=url, config=run_config)
                        return json.dumps(result.extracted_content, indent=2)[:10000] if result.extracted_content else "Error: No structured content extracted"
                else:
                    return f"Error: Unknown action '{action}'. Use scrape, deep_crawl, or extract_structured"

        return asyncio.run(_crawl())
    except ImportError as e:
        return f"Error: crawl4ai not installed. Run: pip install crawl4ai. Details: {e}"
    except Exception as e:
        return f"Crawl4ai error: {str(e)}"


def firecrawl_tool(action: str = "scrape", url: str = "", query: str = "",
                   limit: int = 10, formats: str = "markdown", mode: str = "", **kwargs) -> str:
    try:
        from firecrawl import Firecrawl
        from settings_manager import settings_manager as _sm
        limit = int(limit) if limit else 10

        if not mode:
            mode = _sm.get("firecrawl_mode", "self_hosted")

        if mode == "self_hosted":
            base_url = _sm.get("firecrawl_url", "http://localhost:3001")
            app = Firecrawl(api_url=base_url)
        else:
            api_key = _sm.get("firecrawl_api_key", "")
            if not api_key:
                return "Error: Firecrawl API key not configured. Set firecrawl_api_key in Settings."
            app = Firecrawl(api_key=api_key)

        if action == "scrape":
            if not url:
                return "Error: URL is required for scrape"
            result = app.scrape(url, formats=[formats])
            return result.markdown if formats == "markdown" else result.html
        elif action == "crawl":
            if not url:
                return "Error: URL is required for crawl"
            result = app.crawl(url, limit=limit, scrape_options={"formats": [formats]})
            return "\n\n---\n\n".join([d.markdown for d in result.data]) if result.data else "Error: No data returned"
        elif action == "map":
            if not url:
                return "Error: URL is required for map"
            result = app.map(url)
            return json.dumps(result.links, indent=2)[:10000]
        elif action == "search":
            if not query:
                return "Error: Query is required for search"
            result = app.search(query, limit=limit)
            return json.dumps([{"title": r.title, "url": r.url, "markdown": r.markdown} for r in result.data.web], indent=2)[:10000]
        elif action == "agent":
            if not query:
                return "Error: Query is required for agent"
            result = app.agent(prompt=query)
            return result.data.result
        else:
            return f"Error: Unknown action '{action}'. Use scrape, crawl, map, search, or agent"
    except ImportError as e:
        return f"Error: firecrawl-py not installed. Run: pip install firecrawl-py. Details: {e}"
    except Exception as e:
        return f"Firecrawl error: {str(e)}"


def crawlee_tool(action: str = "scrape_urls", urls: str = "", max_requests: int = 10,
                 crawler_type: str = "playwright", proxy_url: str = "",
                 javascript_code: str = "", **kwargs) -> str:
    try:
        import asyncio
        from crawlee.crawlers import PlaywrightCrawler, BeautifulSoupCrawler
        max_requests = int(max_requests) if max_requests else 10
        url_list = [u.strip() for u in urls.split(",") if u.strip()]

        if not url_list:
            return "Error: No URLs provided"

        async def _crawlee_playwright():
            crawler = PlaywrightCrawler(max_requests_per_crawl=max_requests)

            @crawler.router.default_handler
            async def handler(context):
                data = {
                    "url": context.request.url,
                    "title": await context.page.title(),
                    "content": await context.page.content(),
                }
                await context.push_data(data)

            result = await crawler.run(url_list)
            return json.dumps(result.items, indent=2)[:10000]

        async def _crawlee_beautifulsoup():
            crawler = BeautifulSoupCrawler(max_requests_per_crawl=max_requests)

            @crawler.router.default_handler
            async def handler(context):
                data = {
                    "url": context.request.url,
                    "title": context.soup.title.string if context.soup.title else "",
                    "content": context.soup.get_text()[:5000],
                }
                await context.push_data(data)

            result = await crawler.run(url_list)
            return json.dumps(result.items, indent=2)[:10000]

        if crawler_type == "beautifulsoup":
            return asyncio.run(_crawlee_beautifulsoup())
        return asyncio.run(_crawlee_playwright())
    except ImportError as e:
        return f"Error: crawlee not installed. Run: pip install crawlee[all]. Details: {e}"
    except Exception as e:
        return f"Crawlee error: {str(e)}"


def markitdown_convert(file_path: str = "", input_type: str = "local_file", url: str = "",
                       use_llm: bool = False, llm_model: str = "", **kwargs) -> str:
    try:
        from markitdown import MarkItDown
        kwargs_markitdown = {}
        if use_llm and llm_model:
            from openai import OpenAI
            kwargs_markitdown["llm_client"] = OpenAI()
            kwargs_markitdown["llm_model"] = llm_model

        md = MarkItDown(**kwargs_markitdown)

        if input_type == "url" and url:
            result = md.convert_url(url)
        elif file_path:
            path = Path(file_path)
            if not path.is_absolute():
                path = Path(WORKSPACE_DIR) / file_path
            if not path.exists():
                return f"Error: File not found: {file_path}"
            result = md.convert(str(path))
        else:
            return "Error: Provide file_path or url"

        return result.text_content[:10000]
    except ImportError as e:
        return f"Error: markitdown not installed. Run: pip install markitdown[all]. Details: {e}"
    except Exception as e:
        return f"MarkItDown error: {str(e)}"


TOOL_DEFINITIONS = {
    "web_search": {
        "name": "web_search",
        "description": "Search the web and scrape top results for clean content (Firecrawl-style)",
        "handler": web_search,
        "params": [
            {"name": "query", "type": "string", "required": True, "label": "Search Query"},
            {"name": "num_results", "type": "number", "required": False, "label": "Num Results", "default": 5},
            {"name": "scrape", "type": "boolean", "required": False, "label": "Scrape URLs for full content", "default": True}
        ]
    },
    "web_scraper": {
        "name": "web_scraper",
        "description": "Scrape a webpage and return clean text (strips HTML, scripts, styles — Firecrawl clone)",
        "handler": web_scraper,
        "params": [
            {"name": "url", "type": "string", "required": True, "label": "URL to Scrape"},
            {"name": "max_chars", "type": "number", "required": False, "label": "Max Characters", "default": 8000}
        ]
    },
    "read_file": {
        "name": "read_file",
        "description": "Read the contents of a file (txt, pdf, docx, pptx, csv, xlsx)",
        "handler": read_file,
        "params": [
            {"name": "file_path", "type": "string", "required": True, "label": "File Path"}
        ]
    },
    "write_file": {
        "name": "write_file",
        "description": "Write content to a file",
        "handler": write_file,
        "params": [
            {"name": "file_path", "type": "string", "required": True, "label": "File Path"},
            {"name": "content", "type": "string", "required": True, "label": "Content"}
        ]
    },
    "run_code": {
        "name": "run_code",
        "description": "Execute Python code and return output",
        "handler": run_code,
        "params": [
            {"name": "code", "type": "string", "required": True, "label": "Python Code"}
        ]
    },
    "send_email": {
        "name": "send_email",
        "description": "Send an email via SMTP",
        "handler": send_email,
        "params": [
            {"name": "to", "type": "string", "required": True, "label": "To Email"},
            {"name": "subject", "type": "string", "required": True, "label": "Subject"},
            {"name": "body", "type": "string", "required": True, "label": "Body"},
            {"name": "smtp_server", "type": "string", "required": False, "label": "SMTP Server", "default": "smtp.gmail.com"},
            {"name": "smtp_port", "type": "number", "required": False, "label": "SMTP Port", "default": 587},
            {"name": "username", "type": "string", "required": False, "label": "Username"},
            {"name": "password", "type": "string", "required": False, "label": "Password"}
        ]
    },
    "http_request": {
        "name": "http_request",
        "description": "Make an HTTP request to any REST API",
        "handler": http_request,
        "params": [
            {"name": "url", "type": "string", "required": True, "label": "URL"},
            {"name": "method", "type": "select", "required": True, "label": "Method",
             "options": ["GET", "POST", "PUT", "DELETE", "PATCH"]},
            {"name": "headers", "type": "string", "required": False, "label": "Headers (JSON)"},
            {"name": "body", "type": "string", "required": False, "label": "Body"},
            {"name": "timeout", "type": "number", "required": False, "label": "Timeout (s)", "default": 30}
        ]
    },
    "calculate": {
        "name": "calculate",
        "description": "Evaluate a mathematical expression",
        "handler": calculate,
        "params": [
            {"name": "expression", "type": "string", "required": True, "label": "Expression"}
        ]
    },
    "get_datetime": {
        "name": "get_datetime",
        "description": "Get the current date and time",
        "handler": get_datetime,
        "params": [
            {"name": "format_str", "type": "string", "required": False, "label": "Format",
             "default": "%Y-%m-%d %H:%M:%S"}
        ]
    },
    "run_command": {
        "name": "run_command",
        "description": "Execute a terminal/shell command on the user's PC",
        "handler": run_command,
        "params": [
            {"name": "command", "type": "string", "required": True, "label": "Command"},
            {"name": "working_directory", "type": "string", "required": False, "label": "Working Directory"},
            {"name": "timeout", "type": "number", "required": False, "label": "Timeout (s)", "default": 60}
        ]
    },
    "playwright_browser": {
        "name": "playwright_browser",
        "description": "Browse the web using Playwright (Chromium, Chrome, Brave, Edge, Firefox, WebKit)",
        "handler": playwright_browser,
        "params": [
            {"name": "action", "type": "select", "required": True, "label": "Action",
             "options": ["goto", "click", "type", "extract", "screenshot", "evaluate"]},
            {"name": "url", "type": "string", "required": False, "label": "URL"},
            {"name": "browser", "type": "select", "required": False, "label": "Browser", "default": "chrome",
             "options": ["chrome", "brave", "msedge", "chromium", "firefox", "webkit"]},
            {"name": "selector", "type": "string", "required": False, "label": "CSS Selector"},
            {"name": "text", "type": "string", "required": False, "label": "Text / JS Code"},
            {"name": "screenshot", "type": "string", "required": False, "label": "Screenshot Path"},
            {"name": "wait_seconds", "type": "number", "required": False, "label": "Wait (s)", "default": 3}
        ]
    },
    "list_directory": {
        "name": "list_directory",
        "description": "List files and folders in a directory with sizes",
        "handler": list_directory,
        "params": [
            {"name": "path", "type": "string", "required": False, "label": "Directory Path"},
            {"name": "pattern", "type": "string", "required": False, "label": "Glob Pattern", "default": ""}
        ]
    },
    "search_files": {
        "name": "search_files",
        "description": "Search for files by name pattern or grep content across files",
        "handler": search_files,
        "params": [
            {"name": "query", "type": "string", "required": True, "label": "Search Query"},
            {"name": "path", "type": "string", "required": False, "label": "Search Directory"},
            {"name": "search_type", "type": "select", "required": False, "label": "Search Type",
             "options": ["name", "content"], "default": "name"}
        ]
    },
    "json_query": {
        "name": "json_query",
        "description": "Extract values from JSON using dot-notation path (e.g. data.users.0.name)",
        "handler": json_query,
        "params": [
            {"name": "data", "type": "textarea", "required": True, "label": "JSON Data"},
            {"name": "query", "type": "string", "required": True, "label": "Query Path"}
        ]
    },
    "csv_analyze": {
        "name": "csv_analyze",
        "description": "Analyze CSV files - summary stats, filter, sort, group by columns",
        "handler": csv_analyze,
        "params": [
            {"name": "file_path", "type": "string", "required": True, "label": "CSV File Path"},
            {"name": "operation", "type": "select", "required": True, "label": "Operation",
             "options": ["summary", "head", "filter", "sort", "group"], "default": "summary"},
            {"name": "column", "type": "string", "required": False, "label": "Column Name"},
            {"name": "filter_value", "type": "string", "required": False, "label": "Filter / Sort Direction"},
            {"name": "limit", "type": "number", "required": False, "label": "Row Limit", "default": 20}
        ]
    },
    "database_query": {
        "name": "database_query",
        "description": "Execute SQL queries against a SQLite database",
        "handler": database_query,
        "params": [
            {"name": "db_path", "type": "string", "required": True, "label": "Database Path"},
            {"name": "query", "type": "textarea", "required": True, "label": "SQL Query"},
            {"name": "params", "type": "string", "required": False, "label": "Bind Parameters (JSON array)"}
        ]
    },
    "random_generate": {
        "name": "random_generate",
        "description": "Generate random UUIDs, passwords, numbers, or strings",
        "handler": random_generate,
        "params": [
            {"name": "type", "type": "select", "required": True, "label": "Type",
             "options": ["uuid", "password", "number", "string"], "default": "uuid"},
            {"name": "length", "type": "number", "required": False, "label": "Length", "default": 16},
            {"name": "count", "type": "number", "required": False, "label": "Count", "default": 1}
        ]
    },
    "diff_texts": {
        "name": "diff_texts",
        "description": "Compare two texts and show unified diff",
        "handler": diff_texts,
        "params": [
            {"name": "text_a", "type": "textarea", "required": True, "label": "Text A"},
            {"name": "text_b", "type": "textarea", "required": True, "label": "Text B"},
            {"name": "context_lines", "type": "number", "required": False, "label": "Context Lines", "default": 3}
        ]
    },
    "hash_data": {
        "name": "hash_data",
        "description": "Generate SHA256, SHA1, or MD5 hash of data",
        "handler": hash_data,
        "params": [
            {"name": "data", "type": "string", "required": True, "label": "Data to Hash"},
            {"name": "algorithm", "type": "select", "required": False, "label": "Algorithm",
             "options": ["sha256", "sha1", "md5"], "default": "sha256"}
        ]
    },
    "rss_read": {
        "name": "rss_read",
        "description": "Read and parse RSS/Atom feeds",
        "handler": rss_read,
        "params": [
            {"name": "url", "type": "string", "required": True, "label": "Feed URL"},
            {"name": "max_items", "type": "number", "required": False, "label": "Max Items", "default": 10}
        ]
    },
    "slack_webhook": {
        "name": "slack_webhook",
        "description": "Send a message to Slack via incoming webhook",
        "handler": slack_webhook,
        "params": [
            {"name": "webhook_url", "type": "string", "required": True, "label": "Webhook URL"},
            {"name": "message", "type": "string", "required": True, "label": "Message"},
            {"name": "channel", "type": "string", "required": False, "label": "Channel Override"},
            {"name": "username", "type": "string", "required": False, "label": "Bot Username", "default": "OllamaFlow"}
        ]
    },
    "web_research": {
        "name": "web_research",
        "description": "Search the web and scrape pages using Playwright + Brave browser to bypass anti-bot defenses (403/429 errors)",
        "handler": web_research,
        "params": [
            {"name": "action", "type": "select", "required": True, "label": "Action",
             "options": ["search", "search_and_scrape", "goto", "extract", "screenshot", "click", "type", "evaluate"]},
            {"name": "query", "type": "string", "required": False, "label": "Search Query / JS Code"},
            {"name": "url", "type": "string", "required": False, "label": "URL"},
            {"name": "engine", "type": "select", "required": False, "label": "Search Engine", "default": "google",
             "options": ["google", "bing", "duckduckgo", "brave"]},
            {"name": "max_results", "type": "number", "required": False, "label": "Max Results", "default": 5},
            {"name": "selector", "type": "string", "required": False, "label": "CSS Selector (for extract/click/type)"},
            {"name": "screenshot", "type": "string", "required": False, "label": "Screenshot Path"},
            {"name": "wait_seconds", "type": "number", "required": False, "label": "Wait (s)", "default": 3}
        ]
    },
    "youtube_transcript": {
        "name": "youtube_transcript",
        "description": "Extract transcript/subtitles from a YouTube video",
        "handler": youtube_transcript,
        "params": [
            {"name": "video_url", "type": "string", "required": True, "label": "YouTube Video URL"},
            {"name": "language", "type": "string", "required": False, "label": "Language Code", "default": "en"}
        ]
    },
    "clipboard_copy": {
        "name": "clipboard_copy",
        "description": "Copy text to the system clipboard",
        "handler": clipboard_copy,
        "params": [
            {"name": "text", "type": "string", "required": True, "label": "Text to Copy"}
        ]
    },
    "clipboard_paste": {
        "name": "clipboard_paste",
        "description": "Read text from the system clipboard",
        "handler": clipboard_paste,
        "params": []
    },
    "rate_limiter": {
        "name": "rate_limiter",
        "description": "Rate-limit API calls with sliding window (in-memory)",
        "handler": rate_limiter,
        "params": [
            {"name": "max_requests", "type": "number", "required": False, "label": "Max Requests", "default": 10},
            {"name": "window_seconds", "type": "number", "required": False, "label": "Window (seconds)", "default": 60},
            {"name": "bucket", "type": "string", "required": False, "label": "Bucket Name", "default": "default"}
        ]
    },
    "file_watcher": {
        "name": "file_watcher",
        "description": "List recently modified files in a directory",
        "handler": file_watcher,
        "params": [
            {"name": "path", "type": "string", "required": False, "label": "Directory Path"},
            {"name": "pattern", "type": "string", "required": False, "label": "Glob Pattern", "default": "*"},
            {"name": "action", "type": "select", "required": False, "label": "Action", "options": ["list_changes"], "default": "list_changes"},
            {"name": "since_minutes", "type": "number", "required": False, "label": "Since (minutes)", "default": 60}
        ]
    },
    "browser_use": {
        "name": "browser_use",
        "description": "AI-driven browser agent — give it a natural language task and it automates the web autonomously",
        "handler": browser_use_tool,
        "params": [
            {"name": "task", "type": "textarea", "required": True, "label": "Task Description"},
            {"name": "llm_provider", "type": "select", "required": True, "label": "LLM Provider",
             "options": ["ollama", "openai", "anthropic"]},
            {"name": "llm_model", "type": "string", "required": False, "label": "Model (blank = default)"},
            {"name": "headless", "type": "boolean", "required": False, "label": "Headless Mode", "default": True},
            {"name": "allowed_domains", "type": "string", "required": False, "label": "Allowed Domains (comma-separated)"},
            {"name": "max_steps", "type": "number", "required": False, "label": "Max Steps", "default": 25}
        ]
    },
    "crawl4ai": {
        "name": "crawl4ai",
        "description": "LLM-friendly web crawler — returns clean Markdown optimized for AI. Supports deep crawling and anti-bot detection",
        "handler": crawl4ai_tool,
        "params": [
            {"name": "action", "type": "select", "required": True, "label": "Action",
             "options": ["scrape", "deep_crawl", "extract_structured"]},
            {"name": "url", "type": "string", "required": True, "label": "URL"},
            {"name": "max_pages", "type": "number", "required": False, "label": "Max Pages (deep crawl)", "default": 5},
            {"name": "css_selector", "type": "string", "required": False, "label": "CSS Selector"},
            {"name": "javascript", "type": "textarea", "required": False, "label": "JavaScript to Execute"},
            {"name": "cache", "type": "boolean", "required": False, "label": "Enable Caching", "default": True},
            {"name": "fit_markdown", "type": "boolean", "required": False, "label": "Filter Noise (fit_markdown)", "default": True}
        ]
    },
    "firecrawl": {
        "name": "firecrawl",
        "description": "Web scraping API — scrape, crawl, map, search, or use AI agent. Supports self-hosted Docker and cloud",
        "handler": firecrawl_tool,
        "params": [
            {"name": "action", "type": "select", "required": True, "label": "Action",
             "options": ["scrape", "crawl", "map", "search", "agent"]},
            {"name": "url", "type": "string", "required": False, "label": "URL"},
            {"name": "query", "type": "string", "required": False, "label": "Search Query / Agent Prompt"},
            {"name": "limit", "type": "number", "required": False, "label": "Limit", "default": 10},
            {"name": "formats", "type": "select", "required": False, "label": "Output Format",
             "options": ["markdown", "html"], "default": "markdown"},
            {"name": "mode", "type": "select", "required": False, "label": "Mode",
             "options": ["self_hosted", "cloud"], "default": "self_hosted"}
        ]
    },
    "crawlee": {
        "name": "crawlee",
        "description": "Anti-bot web scraping with proxy rotation, fingerprint spoofing, and automatic retries. Appears human-like",
        "handler": crawlee_tool,
        "params": [
            {"name": "action", "type": "select", "required": True, "label": "Action",
             "options": ["scrape_urls", "deep_crawl"]},
            {"name": "urls", "type": "string", "required": True, "label": "URLs (comma-separated)"},
            {"name": "max_requests", "type": "number", "required": False, "label": "Max Requests", "default": 10},
            {"name": "crawler_type", "type": "select", "required": False, "label": "Crawler Type",
             "options": ["playwright", "beautifulsoup"], "default": "playwright"},
            {"name": "proxy_url", "type": "string", "required": False, "label": "Proxy URL"},
            {"name": "javascript_code", "type": "textarea", "required": False, "label": "JavaScript to Execute"}
        ]
    },
    "markitdown": {
        "name": "markitdown",
        "description": "Convert files to Markdown — PDF, Word, PowerPoint, Excel, images (OCR), audio (transcription), HTML, YouTube",
        "handler": markitdown_convert,
        "params": [
            {"name": "file_path", "type": "string", "required": False, "label": "File Path"},
            {"name": "input_type", "type": "select", "required": False, "label": "Input Type",
             "options": ["local_file", "url"], "default": "local_file"},
            {"name": "url", "type": "string", "required": False, "label": "URL (for url input type)"},
            {"name": "use_llm", "type": "boolean", "required": False, "label": "Use LLM for Image Descriptions", "default": False},
            {"name": "llm_model", "type": "string", "required": False, "label": "LLM Model (for images)"}
        ]
    }
}


def execute_tool(tool_name: str, params: Dict[str, Any]) -> str:
    if tool_name not in TOOL_DEFINITIONS:
        return f"Error: Unknown tool '{tool_name}'"
    tool = TOOL_DEFINITIONS[tool_name]
    try:
        return tool["handler"](**params)
    except TypeError as e:
        return f"Error: Invalid parameters for {tool_name}: {str(e)}"
    except Exception as e:
        return f"Tool error ({tool_name}): {str(e)}"


def get_tools_list() -> list:
    return [
        {
            "name": t["name"],
            "description": t["description"],
            "params": t["params"]
        }
        for t in TOOL_DEFINITIONS.values()
    ]
